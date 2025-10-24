from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# Create presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Social and Cultural Determinants of Health and Disease: A Comprehensive Analysis"
subtitle.text = "Postgraduate Seminar Presentation\nDepartment of Community Medicine\n[Your Name]\n[Date]"

# Slide 2: Learning Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Session Learning Objectives'
tf = body_shape.text_frame
tf.text = 'Upon completion of this seminar, participants will:'
p = tf.add_paragraph()
p.text = '• Define and differentiate social and cultural determinants of health.'
p = tf.add_paragraph()
p.text = '• Analyze the impact of these determinants on health outcomes using evidence-based data.'
p = tf.add_paragraph()
p.text = '• Evaluate Indian-specific contexts and government interventions.'
p = tf.add_paragraph()
p.text = '• Synthesize findings from high-impact research, including the Lancet Commission, to propose actionable recommendations.'
p = tf.add_paragraph()
p.text = '• Engage in discussions on policy implications for community medicine.'

# Slide 3: Presentation Outline
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Presentation Outline'
tf = body_shape.text_frame
tf.text = '1. Introduction to Determinants'
p = tf.add_paragraph()
p.text = '2. Theoretical Framework'
p = tf.add_paragraph()
p.text = '3. Lancet Commission on Social Determinants of Health (2008)'
p = tf.add_paragraph()
p.text = '4. Social Determinants in Detail'
p = tf.add_paragraph()
p.text = '5. Cultural Determinants in Detail'
p = tf.add_paragraph()
p.text = '6. Indian Socio-Cultural Context'
p = tf.add_paragraph()
p.text = '7. Evidence from Research'
p = tf.add_paragraph()
p.text = '8. Data Visualizations'
p = tf.add_paragraph()
p.text = '9. Case Studies'
p = tf.add_paragraph()
p.text = '10. Conclusions and Implications'
p = tf.add_paragraph()
p.text = '11. Recommendations and Action Plans'
p = tf.add_paragraph()
p.text = '12. References and Further Reading'

# Slide 4: Introduction to Determinants
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Introduction to Determinants'
tf = body_shape.text_frame
tf.text = 'Social determinants of health (SDOH) encompass the conditions in which people are born, grow, live, work, and age (WHO, 2023).'
p = tf.add_paragraph()
p.text = 'Cultural determinants include shared beliefs, values, norms, and practices that influence health behaviors (CDC, 2024).'
p = tf.add_paragraph()
p.text = 'These factors interact to create health inequities, disproportionately affecting vulnerable populations.'
p = tf.add_paragraph()
p.text = 'In India, rapid economic growth coexists with persistent disparities, making this topic critical for public health.'

# Slide 5: Theoretical Framework
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Theoretical Framework'
tf = body_shape.text_frame
tf.text = 'Based on the WHO Commission on Social Determinants of Health (2008), SDOH operate through pathways like material circumstances, psychosocial factors, and behavioral patterns.'
p = tf.add_paragraph()
p.text = 'Cultural determinants are framed by theories of cultural competence and health belief models (Kleinman, 1980).'
p = tf.add_paragraph()
p.text = 'Intersectionality theory highlights how social identities (e.g., caste, gender) compound health risks (Crenshaw, 1989).'

# Slide 6: Lancet Commission on Social Determinants of Health (2008)
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Lancet Commission on Social Determinants of Health (2008)'
tf = body_shape.text_frame
tf.text = 'Established in 2005, the Commission aimed to address health inequities through action on SDOH.'
p = tf.add_paragraph()
p.text = 'Key Findings:'
p = tf.add_paragraph()
p.text = '• Social determinants are responsible for the majority of health inequities within and between countries.'
p = tf.add_paragraph()
p.text = '• The social gradient in health: Lower socioeconomic position correlates with worse health outcomes.'
p = tf.add_paragraph()
p.text = '• Recommendations: Improve daily living conditions, tackle inequitable distribution of power, money, and resources, and measure the problem.'

# Slide 7: Lancet Commission - Three Areas for Action
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Lancet Commission - Three Areas for Action'
tf = body_shape.text_frame
tf.text = '1. Improve daily living conditions: Circumstances of birth, growth, living, working, and aging.'
p = tf.add_paragraph()
p.text = '2. Tackle the inequitable distribution of power, money, and resources: Structural drivers like policies and governance.'
p = tf.add_paragraph()
p.text = '3. Measure and understand the problem: Expand knowledge base, train workforce, raise awareness.'
p = tf.add_paragraph()
p.text = 'The Commission emphasized universal but proportionate action to address inequities.'

# Slide 8: Updates and Follow-ups (2023)
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Updates and Follow-ups (2023)'
tf = body_shape.text_frame
tf.text = 'In 2023, The Lancet published updates on SDOH, highlighting persistent inequities exacerbated by COVID-19.'
p = tf.add_paragraph()
p.text = 'Key Insights:'
p = tf.add_paragraph()
p.text = '• Widening gaps in health outcomes due to social determinants.'
p = tf.add_paragraph()
p.text = '• Need for renewed commitment to equity in post-pandemic recovery.'
p = tf.add_paragraph()
p.text = '• Emphasis on intersectoral action and health equity monitoring.'

# Slide 9: Social Determinants - Economic Factors
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social Determinants: Economic Factors'
tf = body_shape.text_frame
tf.text = 'Poverty limits access to nutritious food, safe housing, and healthcare, leading to higher disease burden (Marmot, 2005).'
p = tf.add_paragraph()
p.text = 'In India, 21.9% live below the poverty line, correlating with higher malnutrition and infectious diseases (NFHS-5, 2021).'
p = tf.add_paragraph()
p.text = 'Unemployment exacerbates mental health issues and chronic conditions.'

# Slide 10: Social Determinants - Education
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social Determinants: Education'
tf = body_shape.text_frame
tf.text = 'Education enhances health literacy, enabling better decision-making and preventive behaviors (Nutbeam, 2000).'
p = tf.add_paragraph()
p.text = 'India\'s literacy rate is 77.7%, with rural-urban gaps affecting maternal and child health (Census 2011).'
p = tf.add_paragraph()
p.text = 'Low education is linked to higher NCD prevalence.'

# Slide 11: Social Determinants - Healthcare Access
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social Determinants: Healthcare Access'
tf = body_shape.text_frame
tf.text = 'Barriers include affordability, availability, and quality, particularly in rural India.'
p = tf.add_paragraph()
p.text = 'Ayushman Bharat aims to cover 50 crore people, but implementation challenges persist (PMJAY, 2018).'

# Slide 12: Social Determinants - Environment
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social Determinants: Environment'
tf = body_shape.text_frame
tf.text = 'Neighborhood factors like pollution and sanitation influence respiratory and infectious diseases.'
p = tf.add_paragraph()
p.text = 'Swachh Bharat Mission has built over 10 crore toilets, reducing open defecation (SBM, 2014).'

# Slide 13: Social Determinants - Social Context
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social Determinants: Social Context'
tf = body_shape.text_frame
tf.text = 'Discrimination based on caste, gender, and class affects resource allocation and health outcomes.'
p = tf.add_paragraph()
p.text = 'In India, Scheduled Castes face higher morbidity due to social exclusion (NFHS-5).'

# Slide 14: Cultural Determinants - Beliefs
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Cultural Determinants: Beliefs and Practices'
tf = body_shape.text_frame
tf.text = 'Cultural beliefs about illness (e.g., karma, supernatural causes) can delay modern treatment (Kleinman, 1980).'
p = tf.add_paragraph()
p.text = 'In India, stigma around mental health and HIV/AIDS hinders care-seeking.'

# Slide 15: Cultural Determinants - Gender Norms
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Cultural Determinants: Gender Norms'
tf = body_shape.text_frame
tf.text = 'Patriarchal norms limit women\'s access to education and healthcare, affecting maternal health.'
p = tf.add_paragraph()
p.text = 'Son preference contributes to sex-selective abortions and female malnutrition (Beti Bachao Beti Padhao, 2015).'

# Slide 16: Cultural Determinants - Caste System
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Cultural Determinants: Caste System'
tf = body_shape.text_frame
tf.text = 'Caste influences occupation and social status, leading to health disparities.'
p = tf.add_paragraph()
p.text = 'Lower castes have higher rates of anemia and undernutrition (NFHS-5).'

# Slide 17: Cultural Determinants - Language
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Cultural Determinants: Language and Communication'
tf = body_shape.text_frame
tf.text = 'India\'s multilingualism poses barriers to health education and services.'
p = tf.add_paragraph()
p.text = 'Localization of health messages is essential for effective communication.'

# Slide 18: Indian Socio-Cultural Context
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Indian Socio-Cultural Context'
tf = body_shape.text_frame
tf.text = 'India\'s diversity in culture, religion, and geography amplifies determinant effects.'
p = tf.add_paragraph()
p.text = 'Urbanization and migration create new health challenges like stress-related disorders.'
p = tf.add_paragraph()
p.text = 'Government programs address these through integrated approaches.'

# Slide 19: National Health Mission (NHM)
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'National Health Mission (NHM)'
tf = body_shape.text_frame
tf.text = 'Launched in 2005, focuses on universal healthcare access.'
p = tf.add_paragraph()
p.text = 'Achievements: Reduced MMR from 254 to 97 per 100,000 live births (NHM Report, 2022).'

# Slide 20: Ayushman Bharat
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Ayushman Bharat'
tf = body_shape.text_frame
tf.text = 'Provides financial protection for secondary and tertiary care.'
p = tf.add_paragraph()
p.text = 'Covers 10.74 crore families, reducing out-of-pocket expenses (PMJAY, 2023).'

# Slide 21: Swachh Bharat Mission
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Swachh Bharat Mission'
tf = body_shape.text_frame
tf.text = 'Promotes sanitation and hygiene to prevent diseases.'
p = tf.add_paragraph()
p.text = 'Resulted in 100% open defecation-free status in many states (SBM, 2023).'

# Slide 22: Evidence from Research - Global
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Evidence from Research - Global'
tf = body_shape.text_frame
tf.text = 'Marmot (2005, The Lancet): Social gradient in health; lower SES linked to higher mortality.'
p = tf.add_paragraph()
p.text = 'WHO (2023): SDOH account for 30-50% of health outcomes in high-income countries.'

# Slide 23: Evidence from Research - India
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Evidence from Research - India'
tf = body_shape.text_frame
tf.text = 'Cowling et al. (2014, Int J Equity Health): State-wise inequities in SDOH.'
p = tf.add_paragraph()
p.text = 'NFHS-5 (2021): 35% children under 5 stunted, higher in rural and low-SES groups.'

# Slide 24: Key Research Articles
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Key Research Articles'
tf = body_shape.text_frame
tf.text = '1. Marmot (2005). The Lancet. DOI: 10.1016/S0140-6736(05)17728-8'
p = tf.add_paragraph()
p.text = '2. Kleinman (1980). Patients and Healers in the Context of Culture.'
p = tf.add_paragraph()
p.text = '3. Crenshaw (1989). Demarginalizing the Intersection of Race and Sex.'
p = tf.add_paragraph()
p.text = '4. Victora et al. (2021). The Lancet. DOI: 10.1016/S0140-6736(21)00394-9'
p = tf.add_paragraph()
p.text = '5. Chaturvedi et al. (2024). Hypertension. DOI: 10.1161/HYPERTENSIONAHA.123.21354'
p = tf.add_paragraph()
p.text = '6. Commission on Social Determinants of Health (2008). The Lancet.'

# Slide 25: Visualizations - Global Disparities
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_shape = slide.shapes.title
title_shape.text = 'Visualizations: Global Health Disparities'
chart_data = ChartData()
chart_data.categories = ['High-Income', 'Low-Income', 'India']
chart_data.add_series('Life Expectancy', (82, 64, 70))
chart_data.add_series('Under-5 Mortality (per 1000)', (5, 40, 28))
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(8), Inches(5), chart_data).chart
chart.has_title = True
chart.chart_title.text_frame.text = 'Global and Indian Health Indicators'

# Slide 26: Visualizations - Indian Inequities
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_shape = slide.shapes.title
title_shape.text = 'Visualizations: Indian Socioeconomic Inequities'
chart_data = ChartData()
chart_data.categories = ['Rural', 'Urban', 'Scheduled Caste', 'General']
chart_data.add_series('Malnutrition Rate (%)', (40, 25, 45, 20))
chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(1), Inches(8), Inches(5), chart_data).chart
chart.has_title = True
chart.chart_title.text_frame.text = 'Malnutrition by Socioeconomic Groups in India'

# Slide 27: Case Studies - Rural India
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Case Studies - Rural India'
tf = body_shape.text_frame
tf.text = 'Poor sanitation in rural areas leads to high diarrhea incidence; Swachh Bharat reduced it by 20% (SBM Report).'
p = tf.add_paragraph()
p.text = 'Low female literacy correlates with higher maternal mortality (NFHS-5).'

# Slide 28: Case Studies - Urban Slums
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Case Studies - Urban Slums'
tf = body_shape.text_frame
tf.text = 'Overcrowding and pollution increase respiratory diseases; air quality index often exceeds 200 in Delhi.'
p = tf.add_paragraph()
p.text = 'Migration-related stress contributes to mental health issues.'

# Slide 29: Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Conclusions'
tf = body_shape.text_frame
tf.text = 'Social and cultural determinants profoundly influence health and disease patterns in India.'
p = tf.add_paragraph()
p.text = 'The Lancet Commission (2008) and updates (2023) highlight the need for urgent action on inequities.'
p = tf.add_paragraph()
p.text = 'Government programs have made strides, but persistent inequities require sustained efforts.'

# Slide 30: Implications for Community Medicine
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Implications for Community Medicine'
tf = body_shape.text_frame
tf.text = 'Community medicine must integrate SDOH in assessments and interventions.'
p = tf.add_paragraph()
p.text = 'Cultural sensitivity is key to effective health promotion.'
p = tf.add_paragraph()
p.text = 'Training in equity-focused practices is essential for future practitioners.'

# Slide 31: Recommendations - Policy
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Recommendations - Policy Level'
tf = body_shape.text_frame
tf.text = 'Adopt Health in All Policies to integrate SDOH across sectors.'
p = tf.add_paragraph()
p.text = 'Strengthen Ayushman Bharat and NHM for universal coverage.'
p = tf.add_paragraph()
p.text = 'Enhance data systems for monitoring inequities.'

# Slide 32: Recommendations - Community
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Recommendations - Community Level'
tf = body_shape.text_frame
tf.text = 'Promote community-led interventions to address cultural barriers.'
p = tf.add_paragraph()
p.text = 'Educate on health literacy and stigma reduction.'
p = tf.add_paragraph()
p.text = 'Engage local leaders for sustainable change.'

# Slide 33: Recommendations - Research
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Recommendations - Research Level'
tf = body_shape.text_frame
tf.text = 'Conduct longitudinal studies on cultural determinants in India.'
p = tf.add_paragraph()
p.text = 'Evaluate program impacts using mixed methods.'
p = tf.add_paragraph()
p.text = 'Foster interdisciplinary research for holistic insights.'

# Slide 34: Action Plans
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Action Plans'
tf = body_shape.text_frame
tf.text = 'Short-term: Scale up sanitation and nutrition programs.'
p = tf.add_paragraph()
p.text = 'Medium-term: Integrate SDOH in medical curricula.'
p = tf.add_paragraph()
p.text = 'Long-term: Achieve SDG 3 through equity-focused policies.'

# Slide 35: References
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'References'
tf = body_shape.text_frame
tf.text = '• WHO. (2023). Social determinants of health.'
p = tf.add_paragraph()
p.text = '• CDC. (2024). Social Determinants of Health.'
p = tf.add_paragraph()
p.text = '• Marmot, M. (2005). The Lancet. DOI: 10.1016/S0140-6736(05)17728-8'
p = tf.add_paragraph()
p.text = '• Kleinman, A. (1980). Patients and Healers in the Context of Culture.'
p = tf.add_paragraph()
p.text = '• Crenshaw, K. (1989). University of Chicago Legal Forum.'
p = tf.add_paragraph()
p.text = '• NFHS-5. (2021). National Family Health Survey.'
p = tf.add_paragraph()
p.text = '• NHM Report. (2022). National Health Mission.'
p = tf.add_paragraph()
p.text = '• PMJAY. (2023). Ayushman Bharat.'
p = tf.add_paragraph()
p.text = '• Commission on Social Determinants of Health. (2008). The Lancet.'

# Save the presentation
prs.save('seminar_presentation_lancet.pptx')
print("Presentation with Lancet Commission saved as seminar_presentation_lancet.pptx")
