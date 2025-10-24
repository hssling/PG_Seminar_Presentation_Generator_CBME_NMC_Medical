from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

# Create presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Social and Cultural Determinants of Health and Disease"
subtitle.text = "Postgraduate Seminar Presentation\nCommunity Medicine\n[Your Name]\n[Date]"

# Slide 2: Learning Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Session Learning Objectives'
tf = body_shape.text_frame
tf.text = 'By the end of this seminar, participants will be able to:'
p = tf.add_paragraph()
p.text = '• Understand the definitions and concepts of social and cultural determinants of health and disease.'
p = tf.add_paragraph()
p.text = '• Identify key factors influencing health outcomes based on literature and evidence.'
p = tf.add_paragraph()
p.text = '• Analyze scientific data and information to arrive at logical conclusions.'
p = tf.add_paragraph()
p.text = '• Discuss recommendations for addressing these determinants in public health practice.'

# Slide 3: Outline
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Presentation Outline'
tf = body_shape.text_frame
tf.text = '1. Introduction'
p = tf.add_paragraph()
p.text = '2. Social Determinants of Health'
p = tf.add_paragraph()
p.text = '3. Cultural Determinants of Health'
p = tf.add_paragraph()
p.text = '4. Evidence and Data'
p = tf.add_paragraph()
p.text = '5. Conclusions'
p = tf.add_paragraph()
p.text = '6. Recommendations'
p = tf.add_paragraph()
p.text = '7. References'

# Slide 4: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Introduction'
tf = body_shape.text_frame
tf.text = 'Social and cultural determinants of health refer to the conditions in which people are born, grow, live, work, and age, including access to power, money, and resources (WHO).'
p = tf.add_paragraph()
p.text = 'These determinants influence health inequities, with lower socioeconomic positions linked to worse health outcomes.'
p = tf.add_paragraph()
p.text = 'Cultural factors, such as beliefs and traditions, also shape health behaviors and disease patterns.'
p = tf.add_paragraph()
p.text = 'This presentation explores these determinants, supported by evidence from academic sources, to highlight their impact on health and disease.'

# Slide 5: Social Determinants of Health
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social Determinants of Health'
tf = body_shape.text_frame
tf.text = 'Social determinants include nonmedical factors like economic stability, education, healthcare access, neighborhood environment, and social context (CDC).'
p = tf.add_paragraph()
p.text = 'Examples:'
p = tf.add_paragraph()
p.text = '• Poverty and unemployment increase risk of illness.'
p = tf.add_paragraph()
p.text = '• Education level affects health literacy and preventive behaviors.'
p = tf.add_paragraph()
p.text = '• Safe housing and transportation influence overall well-being.'
p = tf.add_paragraph()
p.text = 'These factors create a social gradient where lower positions correlate with poorer health (WHO).'

# Slide 6: Cultural Determinants of Health
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Cultural Determinants of Health'
tf = body_shape.text_frame
tf.text = 'Cultural determinants encompass beliefs, norms, traditions, and practices that affect health behaviors.'
p = tf.add_paragraph()
p.text = 'Examples:'
p = tf.add_paragraph()
p.text = '• Dietary habits influenced by cultural traditions (e.g., high-fat diets in certain cultures leading to cardiovascular diseases).'
p = tf.add_paragraph()
p.text = '• Health-seeking behaviors: Stigma around mental health in some cultures delays treatment.'
p = tf.add_paragraph()
p.text = '• Language barriers affecting access to healthcare services.'
p = tf.add_paragraph()
p.text = 'These factors interact with social determinants to exacerbate health disparities.'

# Slide 7: Evidence and Data
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Evidence and Data'
tf = body_shape.text_frame
tf.text = 'Key Evidence from WHO and CDC:'
p = tf.add_paragraph()
p.text = '• 18-year difference in life expectancy between high- and low-income countries.'
p = tf.add_paragraph()
p.text = '• Majority of premature NCD deaths occur in low- and middle-income countries.'
p = tf.add_paragraph()
p.text = '• Under-5 mortality rate 8 times higher in Africa than Europe.'
p = tf.add_paragraph()
p.text = '• Poverty strongly correlates with poorer health and higher premature death risk.'
p = tf.add_paragraph()
p.text = 'Studies show social determinants outweigh genetic factors and healthcare access in influencing health outcomes.'

# Slide 8: Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Conclusions'
tf = body_shape.text_frame
tf.text = 'Social and cultural determinants are major drivers of health inequities and disease patterns.'
p = tf.add_paragraph()
p.text = 'Addressing these requires multi-sectoral action beyond healthcare.'
p = tf.add_paragraph()
p.text = 'They create unfair differences in health status that are avoidable with appropriate interventions.'
p = tf.add_paragraph()
p.text = 'Understanding these determinants is essential for promoting health equity in community medicine.'

# Slide 9: Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Recommendations'
tf = body_shape.text_frame
tf.text = '1. Improve daily living conditions through better housing, education, and job opportunities.'
p = tf.add_paragraph()
p.text = '2. Tackle inequitable distribution of resources via policies on taxation and social protection.'
p = tf.add_paragraph()
p.text = '3. Enhance data collection and awareness to measure impacts and train health professionals.'
p = tf.add_paragraph()
p.text = '4. Promote cultural competence in healthcare to address cultural barriers.'
p = tf.add_paragraph()
p.text = '5. Collaborate with government, private sector, and communities for sustainable change.'

# Slide 10: References
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'References'
tf = body_shape.text_frame
tf.text = '• World Health Organization. (2023). Social determinants of health. Retrieved from https://www.who.int/health-topics/social-determinants-of-health'
p = tf.add_paragraph()
p.text = '• Centers for Disease Control and Prevention. (2024). Social Determinants of Health (SDOH). Retrieved from https://www.cdc.gov/socialdeterminants/index.html'
p = tf.add_paragraph()
p.text = '• Additional sources: Academic articles on cultural determinants (e.g., NCBI, Health Affairs).'

# Save the presentation
prs.save('seminar_presentation.pptx')
print("Presentation saved as seminar_presentation.pptx")
