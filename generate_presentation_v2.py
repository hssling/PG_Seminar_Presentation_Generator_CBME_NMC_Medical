from pptx import Presentation

# Create presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Social and Cultural Determinants of Health and Disease"
subtitle.text = "Postgraduate Seminar Presentation\nCommunity Medicine\n[Your Name]\n[Date]"

# Slide 2: Learning Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Session Learning Objectives'
tf = body_shape.text_frame
tf.text = 'By the end of this seminar, participants will be able to:'
p = tf.add_paragraph()
p.text = '• Understand the definitions and concepts of social and cultural determinants of health and disease.'
p = tf.add_paragraph()
p.text = '• Identify key factors influencing health outcomes based on literature and evidence.'
p = tf.add_paragraph()
p.text = '• Analyze scientific data and information to arrive at logical conclusions.'
p = tf.add_paragraph()
p.text = '• Discuss recommendations for addressing these determinants in public health practice.'
p = tf.add_paragraph()
p.text = '• Explore Indian context, including government programs and efforts to mitigate these determinants.'

# Slide 3: Presentation Outline
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Presentation Outline'
tf = body_shape.text_frame
tf.text = '1. Introduction'
p = tf.add_paragraph()
p.text = '2. Social Determinants of Health'
p = tf.add_paragraph()
p.text = '3. Cultural Determinants of Health'
p = tf.add_paragraph()
p.text = '4. Indian Context and Programs'
p = tf.add_paragraph()
p.text = '5. Evidence and Data'
p = tf.add_paragraph()
p.text = '6. Conclusions'
p = tf.add_paragraph()
p.text = '7. Recommendations and Action Plans'
p = tf.add_paragraph()
p.text = '8. References'

# Slide 4: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Introduction'
tf = body_shape.text_frame
tf.text = 'Social determinants of health (SDOH) are the conditions in which people are born, grow, live, work, and age, including access to power, money, and resources (WHO).'
p = tf.add_paragraph()
p.text = 'Cultural determinants include beliefs, norms, and traditions that shape health behaviors.'
p = tf.add_paragraph()
p.text = 'These factors lead to health inequities, with lower socioeconomic positions associated with worse health outcomes.'
p = tf.add_paragraph()
p.text = 'In India, these determinants are influenced by factors like caste, gender, and urbanization, exacerbating disparities.'
p = tf.add_paragraph()
p.text = 'This presentation draws from global and Indian evidence to promote comprehensive understanding and action.'

# Slide 5: Social Determinants - Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social Determinants of Health - Overview'
tf = body_shape.text_frame
tf.text = 'SDOH are nonmedical factors influencing health outcomes (CDC).'
p = tf.add_paragraph()
p.text = 'Key areas: Economic stability, education, healthcare access, neighborhood environment, social context.'
p = tf.add_paragraph()
p.text = 'They create a social gradient: Lower positions lead to poorer health (WHO).'
p = tf.add_paragraph()
p.text = 'In India, poverty, illiteracy, and poor sanitation are major contributors to disease burden.'

# Slide 6: Economic Stability
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Economic Stability'
tf = body_shape.text_frame
tf.text = 'Poverty and unemployment increase illness risk and reduce access to healthcare.'
p = tf.add_paragraph()
p.text = 'Evidence: Poverty correlates with higher premature death rates (CDC).'
p = tf.add_paragraph()
p.text = 'In India: Over 20% of the population lives below the poverty line, leading to malnutrition and infectious diseases (NFHS-5).'

# Slide 7: Education Access and Quality
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Education Access and Quality'
tf = body_shape.text_frame
tf.text = 'Education affects health literacy and preventive behaviors.'
p = tf.add_paragraph()
p.text = 'Low education linked to higher NCD risks (WHO).'
p = tf.add_paragraph()
p.text = 'In India: Literacy rate 77%, but disparities in rural areas; affects maternal and child health (Census 2011).'

# Slide 8: Healthcare Access and Quality
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Healthcare Access and Quality'
tf = body_shape.text_frame
tf.text = 'Barriers like distance, cost, and quality impact health outcomes.'
p = tf.add_paragraph()
p.text = 'In India: Ayushman Bharat aims to provide universal health coverage, but challenges in rural areas persist.'

# Slide 9: Neighborhood and Built Environment
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Neighborhood and Built Environment'
tf = body_shape.text_frame
tf.text = 'Safe housing, transportation, and clean environments promote health.'
p = tf.add_paragraph()
p.text = 'Polluted air and water increase disease risk (CDC).'
p = tf.add_paragraph()
p.text = 'In India: Swachh Bharat Mission addresses sanitation; urban slums face overcrowding and pollution.'

# Slide 10: Social and Community Context
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Social and Community Context'
tf = body_shape.text_frame
tf.text = 'Social support, discrimination, and community safety influence health.'
p = tf.add_paragraph()
p.text = 'Racism and discrimination drive inequities (CDC).'
p = tf.add_paragraph()
p.text = 'In India: Caste and gender discrimination affect access to resources and healthcare.'

# Slide 11: Cultural Determinants - Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Cultural Determinants of Health - Overview'
tf = body_shape.text_frame
tf.text = 'Cultural factors include beliefs, norms, traditions, and practices affecting health behaviors.'
p = tf.add_paragraph()
p.text = 'They interact with social determinants to shape disease patterns.'
p = tf.add_paragraph()
p.text = 'In India: Diverse cultures influence dietary habits, health-seeking, and stigma around diseases.'

# Slide 12: Cultural Beliefs and Traditions
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Cultural Beliefs and Traditions'
tf = body_shape.text_frame
tf.text = 'Beliefs about illness causation (e.g., supernatural) delay modern treatment.'
p = tf.add_paragraph()
p.text = 'Traditional diets may lead to nutritional deficiencies or excesses.'
p = tf.add_paragraph()
p.text = 'In India: Preference for traditional medicine (Ayurveda) alongside allopathy; stigma on mental health.'

# Slide 13: Gender and Cultural Norms
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Gender and Cultural Norms'
tf = body_shape.text_frame
tf.text = 'Gender roles affect access to education, nutrition, and healthcare.'
p = tf.add_paragraph()
p.text = 'In India: Son preference leads to female infanticide and malnutrition; affects maternal health.'

# Slide 14: Caste and Social Hierarchy
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Caste and Social Hierarchy'
tf = body_shape.text_frame
tf.text = 'Caste system influences occupation, education, and health access.'
p = tf.add_paragraph()
p.text = 'Lower castes face discrimination and poorer health outcomes.'
p = tf.add_paragraph()
p.text = 'In India: Scheduled Castes/Tribes have higher disease burden due to social exclusion.'

# Slide 15: Language and Communication
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Language and Communication'
tf = body_shape.text_frame
tf.text = 'Language barriers hinder healthcare access and health education.'
p = tf.add_paragraph()
p.text = 'In India: Multilingual population; health messages need localization.'

# Slide 16: Indian Context - Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Indian Context - Overview'
tf = body_shape.text_frame
tf.text = 'India faces unique challenges: Rapid urbanization, poverty, and cultural diversity.'
p = tf.add_paragraph()
p.text = 'Government initiatives address SDOH through health and social programs.'
p = tf.add_paragraph()
p.text = 'Focus on equity to reduce disparities in health outcomes.'

# Slide 17: National Health Mission (NHM)
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'National Health Mission (NHM)'
tf = body_shape.text_frame
tf.text = 'Launched in 2005, aims to provide universal access to equitable, affordable healthcare.'
p = tf.add_paragraph()
p.text = 'Addresses SDOH through maternal/child health, disease control, and health systems strengthening.'
p = tf.add_paragraph()
p.text = 'Impact: Reduced infant mortality, improved immunization (NHM Reports).'

# Slide 18: Ayushman Bharat
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Ayushman Bharat'
tf = body_shape.text_frame
tf.text = 'Pradhan Mantri Jan Arogya Yojana (PMJAY) provides health insurance to 50 crore beneficiaries.'
p = tf.add_paragraph()
p.text = 'Targets economic barriers to healthcare access.'
p = tf.add_paragraph()
p.text = 'Health and Wellness Centres promote preventive care.'

# Slide 19: Swachh Bharat Mission
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Swachh Bharat Mission'
tf = body_shape.text_frame
tf.text = 'Aims to achieve open defecation-free India and improve sanitation.'
p = tf.add_paragraph()
p.text = 'Addresses neighborhood environment and reduces infectious diseases.'
p = tf.add_paragraph()
p.text = 'Impact: Over 10 crore toilets built, reducing diarrhea incidence (SBM Reports).'

# Slide 20: Other Efforts
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Other Efforts'
tf = body_shape.text_frame
tf.text = 'Mid-Day Meal Scheme: Improves nutrition and education for children.'
p = tf.add_paragraph()
p.text = 'Beti Bachao Beti Padhao: Addresses gender disparities in health and education.'
p = tf.add_paragraph()
p.text = 'NITI Aayog initiatives for health equity.'

# Slide 21: Evidence and Data - Global
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Evidence and Data - Global'
tf = body_shape.text_frame
tf.text = '18-year life expectancy gap between high- and low-income countries (WHO).'
p = tf.add_paragraph()
p.text = 'Under-5 mortality 8 times higher in Africa than Europe.'
p = tf.add_paragraph()
p.text = 'Poverty correlates with higher NCD and infectious disease risks.'

# Slide 22: Evidence and Data - India
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Evidence and Data - India'
tf = body_shape.text_frame
tf.text = 'Life expectancy: 70 years, but disparities by state and socioeconomic status (SRS 2020).'
p = tf.add_paragraph()
p.text = 'Infant mortality: 28 per 1000 live births, higher in rural areas (NFHS-5).'
p = tf.add_paragraph()
p.text = 'Malnutrition affects 35% of children under 5 (NFHS-5).'

# Slide 23: Case Studies
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Case Studies'
tf = body_shape.text_frame
tf.text = 'Rural India: Poor sanitation leads to high diarrhea rates; Swachh Bharat reduced incidence.'
p = tf.add_paragraph()
p.text = 'Urban Slums: Overcrowding and pollution increase respiratory diseases.'
p = tf.add_paragraph()
p.text = 'Tribal Communities: Cultural isolation affects healthcare access.'

# Slide 24: Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Conclusions'
tf = body_shape.text_frame
tf.text = 'Social and cultural determinants are key drivers of health inequities in India and globally.'
p = tf.add_paragraph()
p.text = 'Addressing them requires integrated approaches beyond healthcare.'
p = tf.add_paragraph()
p.text = 'Indian programs like NHM and Ayushman Bharat demonstrate progress, but challenges remain.'
p = tf.add_paragraph()
p.text = 'Comprehensive action can lead to healthier populations and reduced disparities.'

# Slide 25: Recommendations - Policy Level
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Recommendations - Policy Level'
tf = body_shape.text_frame
tf.text = 'Strengthen universal health coverage through Ayushman Bharat expansion.'
p = tf.add_paragraph()
p.text = 'Integrate SDOH in all policies (Health in All Policies).'
p = tf.add_paragraph()
p.text = 'Enhance data collection on inequities for targeted interventions.'

# Slide 26: Recommendations - Community Level
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Recommendations - Community Level'
tf = body_shape.text_frame
tf.text = 'Promote cultural competence in healthcare delivery.'
p = tf.add_paragraph()
p.text = 'Community education on health literacy and preventive behaviors.'
p = tf.add_paragraph()
p.text = 'Engage local leaders to address stigma and discrimination.'

# Slide 27: Recommendations - Research and Monitoring
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Recommendations - Research and Monitoring'
tf = body_shape.text_frame
tf.text = 'Conduct studies on cultural determinants in Indian contexts.'
p = tf.add_paragraph()
p.text = 'Monitor program impacts using indicators like NFHS.'
p = tf.add_paragraph()
p.text = 'Foster intersectoral collaboration for sustainable change.'

# Slide 28: Action Plans
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'Action Plans'
tf = body_shape.text_frame
tf.text = 'Short-term: Scale up existing programs like Swachh Bharat.'
p = tf.add_paragraph()
p.text = 'Medium-term: Integrate SDOH in medical education and training.'
p = tf.add_paragraph()
p.text = 'Long-term: Achieve health equity through policy reforms and community empowerment.'

# Slide 29: References
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = 'References'
tf = body_shape.text_frame
tf.text = '• WHO. Social determinants of health. https://www.who.int/health-topics/social-determinants-of-health'
p = tf.add_paragraph()
p.text = '• CDC. Social Determinants of Health. https://www.cdc.gov/socialdeterminants/index.html'
p = tf.add_paragraph()
p.text = '• NHM. National Health Mission. https://nhm.gov.in/'
p = tf.add_paragraph()
p.text = '• NFHS-5. National Family Health Survey. https://rchiips.org/nfhs/'
p = tf.add_paragraph()
p.text = '• Additional: Census of India, SRS Reports.'

# Save the presentation
prs.save('seminar_presentation_detailed.pptx')
print("Detailed presentation saved as seminar_presentation_detailed.pptx")
